#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_deck.py — トーキャリ:AI活用の全体像 11スライド(16:9・日本語ゴシック・統一トンマナ)。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

FONT = "Hiragino Sans GB"
NAVY = RGBColor(0x15, 0x25, 0x40)
TEAL = RGBColor(0x10, 0xB3, 0xA4)
BLUE = RGBColor(0x3A, 0x7B, 0xD5)
INK = RGBColor(0x1B, 0x26, 0x34)
GRAY = RGBColor(0x5A, 0x66, 0x75)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENTS = [TEAL, BLUE, RGBColor(0xE8, 0x8B, 0x2A), RGBColor(0x8A, 0x5B, 0xD6)]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set(tf_para, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tf_para.text = text
    tf_para.alignment = align
    for r in tf_para.runs:
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color


def box(slide, x, y, w, h, fill=None, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def base(slide, title, num, accent):
    box(slide, 0, 0, SW, SH, fill=WHITE)
    box(slide, 0, 0, Inches(0.18), SH, fill=accent)           # left accent bar
    # title
    tf = textbox(slide, Inches(0.6), Inches(0.34), Inches(11.4), Inches(0.9))
    _set(tf.paragraphs[0], title, 27, NAVY, bold=True)
    box(slide, Inches(0.62), Inches(1.18), Inches(2.2), Pt(3), fill=accent)  # underline
    # footer
    ft = textbox(slide, Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.4))
    _set(ft.paragraphs[0], "株式会社OSD ／ トーキャリ", 9, GRAY)
    pn = textbox(slide, Inches(11.8), Inches(7.02), Inches(1.3), Inches(0.4))
    _set(pn.paragraphs[0], f"{num:02d} / 11", 9, GRAY, align=PP_ALIGN.RIGHT)


def bullets(slide, items, x, y, w, h, size=15, gap=8, color=INK):
    tf = textbox(slide, x, y, w, h)
    for i, (txt, sub) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set(p, "● " + txt, size, color, bold=True)
        p.space_after = Pt(2 if sub else gap)
        if sub:
            s = tf.add_paragraph(); _set(s, "    " + sub, size-3, GRAY)
            s.space_after = Pt(gap)


def cards(slide, items, y, h, cols=None, top_color=None):
    cols = cols or len(items)
    margin = Inches(0.6); gaps = Inches(0.25)
    total = SW - margin*2
    cw = Emu(int((total - gaps*(cols-1)) / cols))
    for i, (head, body) in enumerate(items):
        col = i % cols
        x = Emu(int(margin + col*(cw + gaps)))
        ac = ACCENTS[i % len(ACCENTS)] if top_color is None else top_color
        box(slide, x, y, cw, h, fill=LIGHT)
        box(slide, x, y, cw, Inches(0.09), fill=ac)
        tf = textbox(slide, Emu(int(x+Inches(0.18))), Emu(int(y+Inches(0.22))), Emu(int(cw-Inches(0.36))), Emu(int(h-Inches(0.4))))
        _set(tf.paragraphs[0], head, 14, NAVY, bold=True)
        if body:
            b = tf.add_paragraph(); _set(b, body, 11, GRAY); b.space_before = Pt(4)


# ---------- 1. タイトル ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, 0, Inches(5.05), SW, Inches(0.07), fill=TEAL)
tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.6), MSO_ANCHOR.MIDDLE)
_set(tf.paragraphs[0], "トーキャリ：AI活用の全体像", 44, WHITE, bold=True)
tf2 = textbox(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(1.2))
_set(tf2.paragraphs[0], "AIが書き・描き・検品し・直し・デプロイする 自己改善型 量産パイプライン", 20, TEAL, bold=True)
tf3 = textbox(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.8))
_set(tf3.paragraphs[0], "株式会社OSD　2026.06", 16, RGBColor(0xC8, 0xD2, 0xDE))

# ---------- 2. 何を作っているか ----------
s = prs.slides.add_slide(BLANK); base(s, "何を作っているか", 2, TEAL)
tf = textbox(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.9))
_set(tf.paragraphs[0], "就活支援サービス「トーキャリ」── 400社 × 4コンテンツ をAIで量産。配信は LINE / LIFF。", 16, INK, bold=True)
cards(s, [
    ("L1：10コマ漫画", "台本＋画像。企業を10コマで理解"),
    ("L2／L3：動画", "企業紹介・決算分析をAIナレ動画化"),
    ("ルーム", "6人格のAI OB訪問（対話）"),
    ("共通キャラ", "ナナ&ハルキが全社・全コンテンツ横断"),
], Inches(2.7), Inches(2.4), cols=4)
tf2 = textbox(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.8))
_set(tf2.paragraphs[0], "→ 就活生が「企業を知る」入口を、4つの形式で全社ぶん自動生成する。", 14, GRAY)

# ---------- 3. AIモデル ----------
s = prs.slides.add_slide(BLANK); base(s, "使っているAIモデルと役割", 3, BLUE)
cards(s, [
    ("Claude（中核）", "台本生成・戦略/設計・実行(Claude Code)"),
    ("Gemini", "画像生成(キャラ参照で一貫性)＋vision QA自動検品"),
    ("Google TTS", "動画ナレーション音声"),
    ("OpenAI", "キャラクター参照画像"),
], Inches(1.7), Inches(2.3), cols=4)
bullets(s, [
    ("6人格ルーム＝Claude対話。冒頭で「私はAIです」と明示（誠実性）", None),
    ("モデルは適材適所：書く=Claude／描く・検品=Gemini／話す=TTS", None),
], Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.0), size=15)

# ---------- 4. 1社1ハーネス ----------
s = prs.slides.add_slide(BLANK); base(s, 'AIパイプライン＝"1社1ハーネス"', 4, TEAL)
tf = textbox(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.8))
_set(tf.paragraphs[0], "1社につき1本の自動ワークフローを、並列でファンアウト。", 16, INK, bold=True)
flow = ["ファクト収集", "型/フック選定", "台本ドラフト", "lint検証", "画像生成", "QA", "本番投入"]
n = len(flow); margin = Inches(0.6); gap = Inches(0.16)
cw = Emu(int((SW - margin*2 - gap*(n-1)) / n))
y = Inches(3.0); h = Inches(1.1)
for i, step in enumerate(flow):
    x = Emu(int(margin + i*(cw+gap)))
    box(s, x, y, cw, h, fill=LIGHT); box(s, x, y, cw, Inches(0.08), fill=ACCENTS[i % 4])
    t = textbox(s, x, y, cw, h, MSO_ANCHOR.MIDDLE)
    _set(t.paragraphs[0], f"{i+1}\n{step}", 12, NAVY, bold=True, align=PP_ALIGN.CENTER)
tf2 = textbox(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.2))
_set(tf2.paragraphs[0], "人が触るのは「承認ゲート」だけ。残りはAIが自走する。", 17, TEAL, bold=True)

# ---------- 5. 自己改善ループ ----------
s = prs.slides.add_slide(BLANK); base(s, "自己改善する自律ループ（毎時）", 5, BLUE)
loop = ["インターンFB", "拾う", "仕分け", "台本/画像修正", "ゲート通過", "本番反映", "書き戻し"]
n = len(loop); cw = Emu(int((SW - Inches(0.6)*2 - Inches(0.16)*(n-1)) / n)); y = Inches(2.4)
for i, step in enumerate(loop):
    x = Emu(int(Inches(0.6) + i*(cw+Inches(0.16))))
    box(s, x, y, cw, Inches(1.0), fill=NAVY)
    t = textbox(s, x, y, cw, Inches(1.0), MSO_ANCHOR.MIDDLE)
    _set(t.paragraphs[0], step, 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
bullets(s, [
    ("毎時、AIがFBを拾って自動で直す → ゲート通過分だけ本番反映 → シートに書き戻し", None),
    ("明確なバグ＝自動修正／判断が要るもの＝人へエスカレーション", None),
    ("人間がやるのは「承認」と「キルスイッチ」だけ", None),
], Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.4), size=15)

# ---------- 6. 安全・品質ゲート ----------
s = prs.slides.add_slide(BLANK); base(s, "安全・品質ゲート（暴走を物理的に止める）", 6, TEAL)
cards(s, [
    ("決定論的lint", "出典なき数字・採用倍率・汎用表現・キャラ違反・型重複を検出。error台本はSQL非生成＝本番に入らない"),
    ("一般化canary", "対象社"+"以外"+"が不変かをhash検証"),
    ("D1バックアップ", "1コマンドで巻き戻し可"),
    ("Source-or-Silence", "出典なき数字は書かない"),
    ("業界別パイロット", "3社で検証してから残りへ"),
], Inches(1.9), Inches(2.7), cols=5)
tf = textbox(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.9))
_set(tf.paragraphs[0], "→ AIが間違えても、ゲートが本番反映を物理的にブロックする多層防御。", 16, NAVY, bold=True)

# ---------- 7. マルチエージェント ----------
s = prs.slides.add_slide(BLANK); base(s, "マルチエージェント体制", 7, BLUE)
cards(s, [
    ("設計役 Claude（チャット）", "戦略・文書・Notion・レビュー"),
    ("実行役 Claude Code", "コード・デプロイ・レンダリング"),
    ("人間", "承認者（最終ゲート）"),
], Inches(2.3), Inches(2.6), cols=3)
tf = textbox(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.0))
_set(tf.paragraphs[0], "レーン分離で「同時操作による事故」を防止。役割を分けて並行作業。", 16, NAVY, bold=True)

# ---------- 8. インフラ ----------
s = prs.slides.add_slide(BLANK); base(s, "インフラ構成", 8, TEAL)
cards(s, [
    ("Cloudflare", "Workers / D1 / Pages / R2 / Cron"),
    ("GitHub", "コード・CI・バージョン管理"),
    ("GAS", "スプレッドシート・LINE連携"),
    ("LINE / LIFF", "就活生への配信面"),
], Inches(2.3), Inches(2.2), cols=4)
tf = textbox(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.4))
_set(tf.paragraphs[0], "移行中：自律ループを Mac依存 → GitHub Actions（クラウド）へ。", 17, BLUE, bold=True)
p = tf.add_paragraph(); _set(p, "＝完全24時間 無人化（人の在席・端末に依存しない運用）。", 15, GRAY); p.space_before = Pt(6)

# ---------- 9. 知識蓄積 ----------
s = prs.slides.add_slide(BLANK); base(s, "知識が貯まる回路", 9, BLUE)
flow = ["現場の気づき", "蒸留", "知見集に反映", "生成プロンプトが参照", "次の初稿が向上"]
n = len(flow); cw = Emu(int((SW - Inches(0.6)*2 - Inches(0.2)*(n-1)) / n)); y = Inches(2.7)
for i, step in enumerate(flow):
    x = Emu(int(Inches(0.6) + i*(cw+Inches(0.2))))
    box(s, x, y, cw, Inches(1.2), fill=LIGHT); box(s, x, y, cw, Inches(0.09), fill=ACCENTS[i % 4])
    t = textbox(s, x, y, cw, Inches(1.2), MSO_ANCHOR.MIDDLE)
    _set(t.paragraphs[0], step, 13, NAVY, bold=True, align=PP_ALIGN.CENTER)
tf = textbox(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.2))
_set(tf.paragraphs[0], "FBが「横断ルール」として全社に波及。作るほどAIが賢くなる。", 17, TEAL, bold=True)

# ---------- 10. 現在地 ----------
s = prs.slides.add_slide(BLANK); base(s, "現在地", 10, TEAL)
bullets(s, [
    ("ライブ：総合商社10社＋三井物産", None),
    ("投入進行中：金融・コンサル・自動車", None),
    ("台本400社 生成完了 → 画像 → 投入を順次", None),
    ("夜間に人間ゼロで自律デプロイが回る段階に到達", None),
    ("次：クラウド完全無人化／6人格ルームの量産", None),
], Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.8), size=18, gap=14)

# ---------- 11. 一言で / 論点 ----------
s = prs.slides.add_slide(BLANK); base(s, "一言で / 論点", 11, BLUE)
box(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.2), fill=NAVY)
tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.2), MSO_ANCHOR.MIDDLE)
_set(tf.paragraphs[0], "「AIが書き、描き、検品し、直し、デプロイし、人間は承認とキルスイッチだけ」", 19, WHITE, bold=True)
cards(s, [
    ("①自律デプロイの安全担保", "決定論ゲート＋canary＋可逆性"),
    ("②ハルシネーション対策", "Source-or-Silence＋lint"),
    ("③人的負荷の最小化", "毎時ループで自動化"),
    ("④Mac→クラウド完全無人化", "24h運用への移行"),
], Inches(3.2), Inches(2.9), cols=4)

OUT = Path.home() / "projects" / "10koma-shukatsu" / "outputs" / "tokyari_ai_overview.pptx"
prs.save(str(OUT))
print("saved:", OUT)
